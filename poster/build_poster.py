"""
CoachAI exhibition poster builder.

Generates poster.pptx — a single-page A4 portrait poster summarising the
project for the exhibition. Run from this folder:  python3 build_poster.py

Layout (top-to-bottom):
  1. Title bar (navy)        — project name + author + Goldsmiths
  2. Abstract | Introduction — two-column intro
  3. Phone screenshots strip — three mockups with captions
  4. Design + Implementation | Testing & Evaluation — two-column body
  5. Conclusion              — full-width strip with future work
"""

import os
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Colour palette — Goldsmiths-style navy + warm cream accents
# ---------------------------------------------------------------------------

NAVY        = RGBColor(0x1A, 0x2B, 0x4C)  # header background, section headers
NAVY_DEEP   = RGBColor(0x0F, 0x1B, 0x33)  # title-bar accent strip
CREAM       = RGBColor(0xF5, 0xE7, 0xC4)  # subtitle / accent text
ACCENT      = RGBColor(0xE8, 0x9B, 0x3C)  # callouts (orange/gold)
TEXT        = RGBColor(0x1F, 0x29, 0x37)  # main body text
MUTED       = RGBColor(0x4B, 0x55, 0x63)
SOFT_BG     = RGBColor(0xF1, 0xF4, 0xFB)  # subtle card backgrounds
HAIRLINE    = RGBColor(0xCB, 0xD2, 0xDC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
PRIMARY     = RGBColor(0x25, 0x63, 0xEB)  # CoachAI app primary blue (matches frontend)


# ---------------------------------------------------------------------------
# Set up — A4 portrait, single blank slide
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width  = Inches(8.27)
prs.slide_height = Inches(11.69)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_rect(left, top, width, height, fill, line_colour=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_colour is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_colour
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(left, top, width, height, text, *,
             font="Calibri", size=11, bold=False, italic=False,
             colour=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        lines = [text]
    else:
        lines = text

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        font_obj = run.font
        font_obj.name = font
        font_obj.size = Pt(size)
        font_obj.bold = bold
        font_obj.italic = italic
        font_obj.color.rgb = colour
    return tb


def add_section_header(left, top, width, text):
    """A bold dark-navy section title with a thin accent rule under it."""
    add_text(left, top, width, 0.35, text,
             font="Calibri", size=15, bold=True, colour=NAVY,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    # Accent rule
    line = slide.shapes.add_connector(1, Inches(left + 0.02), Inches(top + 0.32),
                                      Inches(left + 0.7), Inches(top + 0.32))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2.25)


def add_bullets(left, top, width, height, items, *,
                size=10, font="Calibri", colour=TEXT, bold_lead=False,
                line_spacing=1.18):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(2.5)

        # Bullet glyph
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.name = font
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = ACCENT

        # If item is (lead, rest): bold the lead phrase
        if isinstance(item, tuple):
            lead, rest = item
            lead_run = p.add_run()
            lead_run.text = lead
            lead_run.font.name = font
            lead_run.font.size = Pt(size)
            lead_run.font.bold = True
            lead_run.font.color.rgb = NAVY
            tail = p.add_run()
            tail.text = " " + rest
            tail.font.name = font
            tail.font.size = Pt(size)
            tail.font.color.rgb = colour
        else:
            run = p.add_run()
            run.text = item
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold_lead
            run.font.color.rgb = colour
    return tb


def add_image(path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 width=Inches(width), height=Inches(height))


# ---------------------------------------------------------------------------
# 1. Title bar  (slimmer than v1 to free vertical space for the Conclusion)
# ---------------------------------------------------------------------------

H_TITLE = 1.05
add_rect(0, 0, 8.27, H_TITLE, NAVY)
# Accent strip along the bottom of the header
add_rect(0, H_TITLE - 0.07, 8.27, 0.07, ACCENT)

add_text(0.4, 0.10, 5.4, 0.55,
         "CoachAI",
         font="Calibri", size=30, bold=True, colour=WHITE,
         align=PP_ALIGN.LEFT, line_spacing=1.0)
add_text(0.4, 0.55, 5.4, 0.45,
         "A Sport-Aware Strength & Conditioning Recommender for Amateur Athletes",
         font="Calibri", size=12.5, italic=True, colour=CREAM,
         align=PP_ALIGN.LEFT, line_spacing=1.05)

# Right-side author info
add_text(5.7, 0.13, 2.4, 0.30,
         "Pawel Sitko",
         font="Calibri", size=13, bold=True, colour=WHITE,
         align=PP_ALIGN.RIGHT, line_spacing=1.0)
add_text(5.7, 0.43, 2.4, 0.26,
         "BSc Computing — Final Year Project",
         font="Calibri", size=9.5, colour=CREAM,
         align=PP_ALIGN.RIGHT, line_spacing=1.0)
add_text(5.7, 0.66, 2.4, 0.26,
         "Goldsmiths, University of London",
         font="Calibri", size=9.5, colour=CREAM,
         align=PP_ALIGN.RIGHT, line_spacing=1.0)


# ---------------------------------------------------------------------------
# 2. Abstract | Introduction & Background
# ---------------------------------------------------------------------------

ROW2_TOP   = H_TITLE + 0.15
ROW2_HT    = 1.55
COL_GAP    = 0.20
COL_W      = (8.27 - 0.8 - COL_GAP) / 2
LEFT_PAD   = 0.40

# Abstract
add_section_header(LEFT_PAD, ROW2_TOP, COL_W, "Abstract")
add_text(LEFT_PAD, ROW2_TOP + 0.40, COL_W, ROW2_HT - 0.40,
         "CoachAI is a mobile-first web app that generates a personalised "
         "five-day strength and conditioning plan from sport, position, goal "
         "and equipment. It supports volleyball and basketball, uses a "
         "transparent rule-based recommender grounded in S&C literature, "
         "optional sign-in for saving plans, and an LLM-augmented beginner "
         "mode for novice users.",
         size=10, colour=TEXT, line_spacing=1.25)

# Introduction
intro_left = LEFT_PAD + COL_W + COL_GAP
add_section_header(intro_left, ROW2_TOP, COL_W, "Introduction & Background")
add_text(intro_left, ROW2_TOP + 0.40, COL_W, ROW2_HT - 0.40,
         "Amateur athletes routinely lack qualified S&C support. Existing AI "
         "fitness apps (FitnessAI, AI Endurance, Fitbit/Gemini) optimise "
         "general training or endurance sports — none address position-"
         "specific demands of indoor team sports. CoachAI fills this gap "
         "with a transparent, evidence-based recommender for sports where "
         "vertical jump and change-of-direction drive performance.",
         size=10, colour=TEXT, line_spacing=1.25)


# ---------------------------------------------------------------------------
# 3. System / phone screenshots strip
# ---------------------------------------------------------------------------

ROW3_TOP = ROW2_TOP + ROW2_HT + 0.15
ROW3_HT  = 2.95

add_section_header(LEFT_PAD, ROW3_TOP, 7.47, "The System")

# Three phone screenshots in a row, sized to fit
SCREENS_DIR = os.path.join(
    os.path.dirname(__file__), "..", "dissertation", "screenshots"
)
phones_top = ROW3_TOP + 0.40
phone_w = 1.40
phone_h = 2.20
phone_gap = 0.18

# Layout: landing | profile-step | generated plan
phone_paths = [
    "IMG_8756.PNG",   # landing
    "IMG_8761.PNG",   # goal selection
    "IMG_8765.PNG",   # generated plan with position note
]
total_phones_w = phone_w * 3 + phone_gap * 2
phones_start_left = LEFT_PAD + (7.47 - total_phones_w) / 2 - 0.05

for i, fn in enumerate(phone_paths):
    add_image(
        os.path.join(SCREENS_DIR, fn),
        phones_start_left + i * (phone_w + phone_gap),
        phones_top,
        phone_w, phone_h,
    )

# Captions under each phone
caption_top = phones_top + phone_h + 0.06
captions = [
    "1. Landing — position, equipment, goal",
    "2. Pick a training goal (Step 2 of 3)",
    "3. Generated plan with position note",
]
for i, cap in enumerate(captions):
    add_text(
        phones_start_left + i * (phone_w + phone_gap), caption_top,
        phone_w, 0.30, cap,
        font="Calibri", size=8.5, italic=True, colour=MUTED,
        align=PP_ALIGN.CENTER, line_spacing=1.1,
    )


# ---------------------------------------------------------------------------
# 4. Implementation | Testing & Evaluation
# ---------------------------------------------------------------------------

ROW4_TOP = ROW3_TOP + ROW3_HT + 0.15
ROW4_HT  = 3.10

# --- Implementation (left) ---
add_section_header(LEFT_PAD, ROW4_TOP, COL_W, "Implementation")
add_bullets(
    LEFT_PAD, ROW4_TOP + 0.42, COL_W, ROW4_HT - 0.50,
    [
        ("Backend.", "Python + Flask. Single REST endpoint generates "
                     "plans from a four-element user profile."),
        ("Recommender.", "Transparent rule-based engine over a curated "
                         "database of 40+ exercises tagged by sport, "
                         "category and equipment level."),
        ("Session templates.", "Five-day weekly structure — lower power, "
                                "upper strength, lower strength, power & "
                                "agility, core — with position notes and "
                                "goal-driven adjustments."),
        ("Cycle 2: accounts.", "Flask-Login + SQLite for optional sign-up, "
                                "per-user plan saves, and persistent "
                                "session-complete tracking."),
        ("Cycle 2: LLM mode.", "Google Gemini rephrases coaching cues for "
                                "novice users at development time; runtime "
                                "serves a static cache — deterministic and "
                                "offline-capable."),
    ],
    size=9.5,
)

# --- Testing & Evaluation (right) ---
eval_left = LEFT_PAD + COL_W + COL_GAP
add_section_header(eval_left, ROW4_TOP, COL_W, "Testing & Evaluation")

# A small stats banner — three big-number callouts
stats_top = ROW4_TOP + 0.45
stat_w = (COL_W - 0.20) / 3
stat_h = 0.95
for i, (big, small) in enumerate([
    ("120/120", "input combinations\nproduced valid plans"),
    ("4.3 / 5", "mean rating —\nexercises feel\nsport-relevant"),
    ("N = 10", "amateur athletes\nin Phase 1 study"),
]):
    sx = eval_left + i * (stat_w + 0.10)
    add_rect(sx, stats_top, stat_w, stat_h, SOFT_BG, line_colour=HAIRLINE)
    add_text(sx, stats_top + 0.08, stat_w, 0.40, big,
             font="Calibri", size=18, bold=True, colour=NAVY,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(sx, stats_top + 0.45, stat_w, stat_h - 0.45, small,
             font="Calibri", size=8, colour=MUTED,
             align=PP_ALIGN.CENTER, line_spacing=1.1)

# Findings text under the stats
findings_top = stats_top + stat_h + 0.12
add_bullets(
    eval_left, findings_top, COL_W, ROW4_HT - (findings_top - ROW4_TOP) - 0.05,
    [
        ("Functional.", "Every supported (sport × position × goal × "
                        "equipment) combination produces a structurally valid "
                        "plan. Met criteria 1–3."),
        ("Sport relevance.", "90% of survey participants rated the plan’s "
                              "exercises as relevant to their sport (≥ 4/5). "
                              "Met criterion 5 in absolute terms."),
        ("Clarity.", "Mean 3.30/5; weaker for beginners (2.25) than advanced "
                     "(4.50). Drove the Cycle 2 beginner-friendly LLM mode."),
        ("Demand for persistence.", "3 of 10 testers asked to save plans / "
                                     "track progress — directly motivated the "
                                     "Cycle 2 account layer."),
    ],
    size=9.5,
)


# ---------------------------------------------------------------------------
# 5. Conclusion & Future Work + footer
# ---------------------------------------------------------------------------

ROW5_TOP = ROW4_TOP + ROW4_HT + 0.15
ROW5_HT  = 1.55

add_section_header(LEFT_PAD, ROW5_TOP, 7.47, "Conclusion & Future Work")

# Two-column body
con_left_w = 4.0
con_right_w = 7.47 - con_left_w - 0.20
add_text(
    LEFT_PAD, ROW5_TOP + 0.42, con_left_w, ROW5_HT - 0.50,
    "An iterative two-cycle design — a rule-based prototype (Assignment 3) "
    "evaluated with amateur athletes, followed by an extended final system "
    "(Assignment 4) addressing the most-cited Phase 1 findings — produced a "
    "transparent, mobile-first S&C recommender that meets six of seven "
    "success criteria. Findings are reported honestly: criterion 4 "
    "(clarity ≥ 4/5) was not met for novice users, motivating the new "
    "beginner-friendly explanations mode.",
    size=10, colour=TEXT, line_spacing=1.22,
)

add_bullets(
    LEFT_PAD + con_left_w + 0.20, ROW5_TOP + 0.42, con_right_w, ROW5_HT - 0.50,
    [
        ("Future work:", "exercise demonstrations / videos for novices."),
        ("", "Wider basketball-specific coverage."),
        ("", "Richer goal-driven differentiation between plans."),
        ("", "Longitudinal study over multiple weeks of training."),
    ],
    size=9.5,
)


# ---------------------------------------------------------------------------
# 6. Footer — Goldsmiths plate
# ---------------------------------------------------------------------------

FOOTER_HT = 0.50
foot_top = 11.69 - FOOTER_HT
add_rect(0, foot_top, 8.27, FOOTER_HT, NAVY_DEEP)
add_text(0.40, foot_top + 0.07, 7.47, 0.36,
         "Goldsmiths, University of London  ·  Department of Computing  ·  "
         "Final Year Project Exhibition 2026",
         font="Calibri", size=9, colour=CREAM,
         align=PP_ALIGN.LEFT, line_spacing=1.1)
add_text(0.40, foot_top + 0.07, 7.47, 0.36,
         "github.com/PawelSitko/CoachAi",
         font="Calibri", size=9, italic=True, colour=CREAM,
         align=PP_ALIGN.RIGHT, line_spacing=1.1)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------

OUT = os.path.join(os.path.dirname(__file__), "poster.pptx")
prs.save(OUT)
print(f"Wrote {OUT}")
